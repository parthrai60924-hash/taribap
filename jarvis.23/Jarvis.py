# ---------------- AUTO INSTALLER ----------------
import subprocess
import sys

packages = [
    "pyttsx3",
    "SpeechRecognition",
    "pyjokes",
    "psutil",
    "openai",
    "python-pptx",
    "python-docx",
    "requests",
    "beautifulsoup4",
    "pywhatkit",
    "pyaudio"
]

def install_packages():
    print("Installing required packages...")
    for p in packages:
        subprocess.call([sys.executable, "-m", "pip", "install", p])
    print("All packages installed")

print("================================")
print("BHARAT JARVIS INSTALLER")
print("1 Install all required packages")
print("2 Start Bharat Jarvis")
print("================================")

choice = input("Enter option: ")

if choice == "1":
    install_packages()
    print("Restart program after installation")
    exit()

# ---------------- IMPORTS ----------------
import webbrowser
import pyttsx3
import speech_recognition as sr
import os
import platform
import datetime
import pyjokes
import difflib
import psutil
from openai import OpenAI
from pptx import Presentation
from docx import Document
import socket
import ctypes
import tkinter as tk
import threading
import math
import requests
from bs4 import BeautifulSoup
import pywhatkit

# ---------------- CONFIG ---------------- #
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
OS_NAME = platform.system()

# ---------------- VISION EFFECT ---------------- #
def start_vision():
    root = tk.Tk()
    root.title("BHARAT JARVIS VISION")
    width = root.winfo_screenwidth()
    height = root.winfo_screenheight()
    root.geometry(f"{width}x{height}")
    root.configure(bg="black")
    canvas = tk.Canvas(root, bg="black", highlightthickness=0)
    canvas.pack(fill="both", expand=True)
    cx = width // 2
    cy = height // 2
    radius = 200
    angle = 0
    def animate():
        nonlocal angle
        canvas.delete("all")
        canvas.create_oval(cx-radius, cy-radius, cx+radius, cy+radius, outline="cyan", width=3)
        canvas.create_oval(cx-120, cy-120, cx+120, cy+120, outline="cyan")
        x = cx + radius * math.cos(math.radians(angle))
        y = cy + radius * math.sin(math.radians(angle))
        canvas.create_line(cx, cy, x, y, fill="cyan", width=3)
        canvas.create_text(cx, cy+260, text="BHARAT JARVIS VISION SYSTEM", fill="cyan", font=("Arial", 20))
        angle += 4
        root.after(40, animate)
    animate()
    root.mainloop()
    # ---------------- GOOGLE SEARCH READER ---------------- #
def google_search_read(query):
    global reading
    reading = True

    try:
        speak("Searching Google")

        # open google tab
        webbrowser.open(f"https://www.google.com/search?q={query}")

        url = f"https://www.google.com/search?q={query}"

        headers = {"User-Agent": "Mozilla/5.0"}

        response = requests.get(url, headers=headers)
        soup = BeautifulSoup(response.text, "html.parser")

        paragraphs = soup.find_all("span")

        for p in paragraphs:

            if not reading:
                engine.stop()
                return

            text = p.get_text().strip()

            if len(text) > 50:

                sentences = text.split(".")

                for sentence in sentences:

                    if not reading:
                        engine.stop()
                        return

                    speak(sentence)

                return

        speak("Sorry I could not find information")

    except Exception as e:
        print(e)
        speak("Google search failed")
# ---------------- OWNER DETAILS ---------------- #
OWNER_NAME = "Parth Rai"
FATHER_NAME = "Sanjay Kumar Rai"
MOTHER_NAME = "Saumay Rai"
SCHOOL = "Sunbeam School Mau"
HELPER = "Amogh Pandey"
AI_NAME = "Bharat Jarvis"

# ---------------- SPEAK ---------------- #
engine = pyttsx3.init('sapi5')
voices = engine.getProperty('voices')
engine.setProperty('voice', voices[0].id)
engine.setProperty('rate', 185)
engine.setProperty('volume', 1.0)
def speak(text):
    print("Bharat Jarvis:", text)
    engine.stop()
    engine.setProperty('volume', 1.0)
    engine.say(str(text))
    engine.runAndWait()

# ---------------- NOTIFICATION ---------------- #
def notify(title, message):
    ctypes.windll.user32.MessageBoxW(0, message, title, 0x40)

# ---------------- VOICE INPUT ---------------- #
def take_command():
    r = sr.Recognizer()
    with sr.Microphone() as source:
        print("Listening...")
        r.adjust_for_ambient_noise(source)
        audio = r.listen(source)
    try:
        command = r.recognize_google(audio)
        print("You said:", command)
        return command.lower()
    except:
        return ""

# ---------------- WORD CORRECTION ---------------- #
def correct_word(command):
    keywords = ["youtube","spotify","wifi","bluetooth","settings",
        "chrome","notepad","word","shutdown","restart",
        "display","time","date","play","close","open",
        "web","app","powerpoint","ppt","make","write",
        "battery","music"]
    words = command.split()
    corrected = []
    for word in words:
        match = difflib.get_close_matches(word, keywords, n=1, cutoff=0.6)
        if match:
            corrected.append(match[0])
        else:
            corrected.append(word)
    return " ".join(corrected)

# ---------------- AI ---------------- #
def ask_ai(question):
    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": question}]
        )
        answer = response.choices[0].message.content
        speak(answer)
    except:
        speak("Sorry, I cannot answer right now.")

# ---------------- SYSTEM FUNCTIONS ---------------- #
def shutdown():
    speak("Shutting down the system.")
    os.system("shutdown /s /t 5")
def restart():
    speak("Restarting the system.")
    os.system("shutdown /r /t 5")
def close_app(name):
    os.system(f"taskkill /im {name}.exe /f")
    speak(f"I closed {name}")
def open_settings():
    os.system("start ms-settings:")
    speak("Opening settings")
def open_wifi():
    os.system("start ms-settings:network-wifi")
    speak("Opening Wi Fi settings")
def open_bluetooth():
    os.system("start ms-settings:bluetooth")
    speak("Opening Bluetooth")
def open_display():
    os.system("start ms-settings:display")
    speak("Opening display settings")

# ---------------- BATTERY ---------------- #
def battery_status():
    battery = psutil.sensors_battery()
    if battery is None:
        speak("I cannot detect battery.")
        return
    percent = battery.percent
    plugged = battery.power_plugged
    speak(f"Battery level is {percent} percent")
    if plugged and percent == 100:
        speak("Battery is fully charged")
    elif not plugged and percent <= 20:
        speak("Battery is low. Please charge")

# ---------------- OFFICE FUNCTIONS ---------------- #
def create_ppt():
    speak("What is the topic")
    topic = take_command()
    if not topic:
        speak("Topic not detected")
        return
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = topic
    slide.placeholders[1].text = "Created by Bharat Jarvis"
    speak("Tell me file name")
    file_name = take_command()
    if not file_name:
        file_name = "presentation"
    prs.save(f"{file_name}.pptx")
    speak("Presentation saved")
def write_notepad():
    speak("File name please")
    file_name = take_command()
    if not file_name:
        file_name = "note"
    speak("Speak content")
    content = take_command()
    with open(f"{file_name}.txt","w") as f:
        f.write(content)
    speak("File saved")
def write_word():
    speak("Document name")
    file_name = take_command()
    if not file_name:
        file_name = "document"
    speak("Speak content")
    content = take_command()
    doc = Document()
    doc.add_paragraph(content)
    doc.save(f"{file_name}.docx")
    speak("Word document saved")

# ---------------- GOOGLE DEFINITION READER (FIXED) ---------------- #
reading = False
def read_google_definition(query):
    global reading
    reading = True
    try:
        speak("Searching Google")
        webbrowser.open(f"https://www.google.com/search?q={query}")
        headers = {"User-Agent":"Mozilla/5.0"}
        response = requests.get(f"https://www.google.com/search?q={query}", headers=headers)
        soup = BeautifulSoup(response.text, "html.parser")
        result = soup.find("div", class_="BNeawe iBp4i AP7Wnd")
        if not result:
            result = soup.find("div", class_="BNeawe s3v9rd AP7Wnd")
        if result:
            text = result.get_text()
            speak("Here is the first definition")
            for part in text.split("."):
                if not reading:
                    engine.stop()
                    return
                speak(part)
        else:
            speak("Sorry, I could not find definition")
    except Exception as e:
        speak("Internet or parsing problem occurred")
        print(e)

# ---------------- MAIN BRAIN ---------------- #
def jarvis(command):
    global reading
    command = command.lower()
    # OWNER INFO
    if "who is the owner of bharat jarvis" in command or "who made you" in command:
        speak(f"The main owner of Bharat Jarvis is {OWNER_NAME}. His father name is {FATHER_NAME}. His mother name is {MOTHER_NAME}. He studies in {SCHOOL}. The helper of this project is {HELPER}.")
    # SHUTDOWN
    elif "shutdown" in command:
        shutdown()
    # RESTART
    elif "restart" in command:
        restart()
    # BATTERY
    elif "battery" in command:
        battery_status()
    # PLAY MUSIC
    elif command.startswith("play"):
        try:
            song = command.replace("play","")
            pywhatkit.playonyt(song)
            speak(f"Playing {song}")
        except:
            speak("Cannot play music")
    # OPEN WEB
    elif command.startswith("open web"):
        text = command.replace("open web","").strip()
        webbrowser.open(f"https://{text}.com")
        speak(f"Opening {text}")
    # OPEN APP
    elif command.startswith("open app"):
        name = command.replace("open app","").strip()
        os.system(f"start {name}")
        speak(f"Opening {name}")
    # SETTINGS
    elif "open settings" in command:
        open_settings()
    elif "open wifi" in command:
        open_wifi()
    elif "open bluetooth" in command:
        open_bluetooth()
    elif "open display" in command:
        open_display()
    # CLOSE
    elif command.startswith("close"):
        name = command.replace("close","").strip()
        close_app(name)
    # TIME
    elif "time" in command:
        now = datetime.datetime.now().strftime("%H:%M")
        speak(f"The time is {now}")
    # DATE
    elif "date" in command:
        today = datetime.datetime.now().strftime("%d %B %Y")
        speak(f"Today is {today}")
    # VISION
    elif "open vision" in command or "start vision" in command:
        speak("Opening Bharat Jarvis vision system")
        vision_thread = threading.Thread(target=start_vision)
        vision_thread.daemon = True
        vision_thread.start()
    # JOKE
    elif "joke" in command:
        speak(pyjokes.get_joke())
    # STOP READING
    elif "stop" in command:
        reading = False
        engine.stop()
        speak("Stopped reading")
    # OFFICE
    elif "make powerpoint" in command:
        create_ppt()
    elif "write note" in command:
        write_notepad()
    elif "write word" in command:
        write_word()
    # GOOGLE SEARCH / DEFINITION
    else:
        read_google_definition(command)

# ---------------- INTERNET CHECK ---------------- #
def is_connected():
    try:
        socket.create_connection(("8.8.8.8",53),timeout=3)
        return True
    except:
        return False

# ---------------- STARTUP ---------------- #
if is_connected():
    notify("BHARAT JARVIS","Sir I am online")
    speak("Sir Bharat Jarvis is online and ready")
else:
    notify("BHARAT JARVIS","Sir I am offline")
    speak("Sir I am offline")

# ---------------- LOOP ---------------- #
while True:
    cmd = take_command()
    if cmd:
        cmd = correct_word(cmd)
        jarvis(cmd)